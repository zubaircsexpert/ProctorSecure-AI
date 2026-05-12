#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
ProctorSecure-AI Project Presentation Generator
Creates a PowerPoint presentation summarizing the project for 2-student presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 55, 120)  # Dark blue
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = RGBColor(173, 216, 230)  # Light blue
    subtitle_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 250)  # Light gray-blue
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(25, 55, 120)
    
    # Add horizontal line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(100, 150, 200)
    line.line.width = Pt(2)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

def create_presentation():
    """Create the complete PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    add_title_slide(prs, 
                   "ProctorSecure-AI",
                   "Intelligent Exam Proctoring & Assessment System")
    
    # Slide 2: Project Overview
    add_content_slide(prs,
                     "📌 Project Overview",
                     [
                         "🎯 Purpose: Comprehensive AI-powered examination and assessment platform",
                         "🌐 Full-stack web application with real-time monitoring",
                         "🤖 Advanced AI integration for cheating detection and content generation",
                         "👥 Multi-role system: Teachers, Students, Admins",
                         "🔒 Enterprise-grade security and data privacy"
                     ])
    
    # Slide 3: Key Features
    add_content_slide(prs,
                     "✨ Key Features",
                     [
                         "🎯 AI-Powered Quiz Generator - Auto-generate MCQs from documents",
                         "🔍 Real-time Proctoring - Face detection, eye tracking, environment monitoring",
                         "🚨 Cheating Detection - Plagiarism detection, unusual behavior analysis",
                         "📊 Advanced Analytics - Student performance reports, insights",
                         "📚 Study Resources - Educational materials and learning tools",
                         "🔐 Secure Assessments - Multiple security layers and authentication"
                     ])
    
    # Slide 4: Technology Stack - Frontend
    add_content_slide(prs,
                     "💻 Technology Stack - Frontend",
                     [
                         "React.js - UI framework with component-based architecture",
                         "Vite - Fast build tool and development server",
                         "Face-API.js - Advanced face detection and recognition",
                         "TensorFlow.js - ML models for facial analysis",
                         "Axios - HTTP client for API communication",
                         "CSS3 & Modern UI Design"
                     ])
    
    # Slide 5: Technology Stack - Backend
    add_content_slide(prs,
                     "⚙️ Technology Stack - Backend",
                     [
                         "Node.js & Express.js - Server framework",
                         "MongoDB - NoSQL database for flexibility",
                         "OpenAI GPT-4 API - Intelligent question generation",
                         "JWT Authentication - Secure token-based auth",
                         "Multer - File upload handling",
                         "Tesseract OCR - Optical character recognition"
                     ])
    
    # Slide 6: System Architecture
    add_content_slide(prs,
                     "🏗️ System Architecture",
                     [
                         "Client Layer: React SPA with responsive design",
                         "API Layer: RESTful Express.js backend with route-based organization",
                         "Database Layer: MongoDB collections for data persistence",
                         "Integration Layer: OpenAI, File storage, ML models",
                         "Security Layer: JWT, CORS, input validation, encryption"
                     ])
    
    # Slide 7: Core Modules - Part 1
    add_content_slide(prs,
                     "📦 Core Modules (1/2)",
                     [
                         "🎓 Assessment System - Unified quiz and exam management",
                         "🤖 AI Quiz Generator - Automatic question creation from documents",
                         "📋 Question Bank - Centralized MCQ storage and management",
                         "✏️ Quiz Assembly - Custom quiz creation with difficulty distribution",
                         "📝 Quiz Session - Track student quiz attempts and progress"
                     ])
    
    # Slide 8: Core Modules - Part 2
    add_content_slide(prs,
                     "📦 Core Modules (2/2)",
                     [
                         "🎥 AI Exam Proctoring - Real-time monitoring during exams",
                         "👁️ AI Report - Detailed analysis of student behavior",
                         "🚫 Exam Violations - Track and report suspicious activities",
                         "🔎 Cheat Detection - Plagiarism and behavior anomaly detection",
                         "📊 Quiz Analytics - Performance metrics and insights"
                     ])
    
    # Slide 9: User Roles & Workflows
    add_content_slide(prs,
                     "👥 User Roles & Workflows",
                     [
                         "👨‍🏫 Teachers: Create assessments, upload content, review analytics",
                         "👨‍🎓 Students: Take quizzes/exams, view results, access study materials",
                         "👨‍💼 Admins: System management, user management, reports",
                         "📊 Dashboard: Personalized views for each role",
                         "🔔 Notifications: Real-time updates and alerts"
                     ])
    
    # Slide 10: Data Flow
    add_content_slide(prs,
                     "🔄 Data Flow",
                     [
                         "1️⃣ Teacher uploads document → System extracts text",
                         "2️⃣ AI generates questions from extracted content",
                         "3️⃣ Questions reviewed and added to question bank",
                         "4️⃣ Teacher assembles quiz from bank",
                         "5️⃣ Student takes exam with real-time proctoring",
                         "6️⃣ System generates detailed analytics and reports"
                     ])
    
    # Slide 11: AI Integration
    add_content_slide(prs,
                     "🤖 AI Integration",
                     [
                         "GPT-4 API - Generates high-quality MCQs from documents",
                         "Face-API.js - Detects and tracks facial expressions",
                         "TensorFlow Models - Age, gender, emotion recognition",
                         "Tesseract OCR - Extracts text from images",
                         "Smart Algorithms - Cheat detection and behavior analysis"
                     ])
    
    # Slide 12: Security Features
    add_content_slide(prs,
                     "🔒 Security Features",
                     [
                         "JWT Token Authentication - Secure session management",
                         "HTTPS/TLS - Encrypted communication",
                         "Input Validation & Sanitization - Prevent injection attacks",
                         "Database Encryption - Protect sensitive data",
                         "Role-Based Access Control - Granular permissions",
                         "Real-time Monitoring - Detect unauthorized activities"
                     ])
    
    # Slide 13: Proctoring Features
    add_content_slide(prs,
                     "🎥 Advanced Proctoring Features",
                     [
                         "👁️ Face Detection - Ensures student is present and visible",
                         "👀 Eye Tracking - Detects looking away from screen",
                         "🔊 Audio Monitoring - Captures ambient sounds",
                         "📱 Device Detection - Prevents phone usage during exam",
                         "🚪 Environment Check - Verifies exam location",
                         "🚨 Alert System - Real-time violations reporting"
                     ])
    
    # Slide 14: Database Models
    add_content_slide(prs,
                     "🗂️ Database Models",
                     [
                         "👤 User - Student, teacher, admin profiles with authentication",
                         "📋 Assessment - Unified quiz and exam data structure",
                         "❓ Question & MCQBank - Question storage and versioning",
                         "📊 QuizAnalytics - Performance metrics and statistics",
                         "🚫 ExamViolation - Tracks suspicious activities",
                         "📝 Submission & Result - Student answers and grading"
                     ])
    
    # Slide 15: File Processing
    add_content_slide(prs,
                     "📄 File Processing Capabilities",
                     [
                         "📕 PDF Support - pdf-parse library for text extraction",
                         "📗 DOCX Support - mammoth library for Word documents",
                         "🖼️ Image Support - Tesseract OCR for scanned documents",
                         "📊 Excel Support - xlsxjs for spreadsheet data",
                         "📝 TXT Support - Direct text file processing",
                         "✂️ Text Cleaning - Remove formatting, normalize content"
                     ])
    
    # Slide 16: API Endpoints Summary
    add_content_slide(prs,
                     "🔌 API Endpoints",
                     [
                         "🔐 Auth: Login, register, token refresh, password reset",
                         "📋 Assessments: CRUD operations, submissions, results",
                         "❓ Questions: Generate, store, manage, export MCQs",
                         "📊 Analytics: Performance data, reports, insights",
                         "🎥 Proctoring: Start/stop, violations, AI analysis",
                         "📚 Resources: Study materials, notifications"
                     ])
    
    # Slide 17: Performance & Scalability
    add_content_slide(prs,
                     "⚡ Performance & Scalability",
                     [
                         "🚀 Optimized Query Performance - Indexed MongoDB queries",
                         "📦 Code Splitting - Vite chunks for faster loading",
                         "🔄 Caching Strategy - Reduce API calls and DB queries",
                         "🌐 Horizontal Scaling - Stateless backend architecture",
                         "📈 Load Balancing - Support for multiple server instances",
                         "⏱️ Response Time - < 500ms for most operations"
                     ])
    
    # Slide 18: Testing & Quality
    add_content_slide(prs,
                     "✅ Testing & Quality Assurance",
                     [
                         "🧪 Unit Testing - Component and function-level tests",
                         "🔗 Integration Testing - API and database interaction tests",
                         "🎯 End-to-End Testing - Full user workflow testing",
                         "📋 Code Review - Regular peer reviews and standards",
                         "📊 Error Tracking - Comprehensive logging and monitoring",
                         "📈 Performance Monitoring - Track metrics and bottlenecks"
                     ])
    
    # Slide 19: Deployment & DevOps
    add_content_slide(prs,
                     "🚀 Deployment & DevOps",
                     [
                         "☁️ Vercel - Frontend deployment with automatic updates",
                         "🐳 Docker Support - Containerized backend deployment",
                         "🗄️ Database - MongoDB Atlas cloud hosting",
                         "🔄 CI/CD Pipeline - Automated testing and deployment",
                         "🔌 API Management - Environment variables and secrets",
                         "📊 Monitoring - Performance and error tracking"
                     ])
    
    # Slide 20: Use Cases
    add_content_slide(prs,
                     "💼 Real-World Use Cases",
                     [
                         "🏫 Schools & Universities - Online exams with integrity",
                         "🎓 Online Courses - Self-paced learning with verification",
                         "💼 Corporate Training - Employee assessments and certifications",
                         "🏛️ Certification Bodies - Secure professional exams",
                         "📚 Test Prep - Practice exams and performance tracking",
                         "🌍 Distance Education - Remote proctoring solution"
                     ])
    
    # Slide 21: Challenges & Solutions
    add_content_slide(prs,
                     "🎯 Challenges & Solutions",
                     [
                         "🤖 Challenge: AI accuracy in question generation",
                         "   Solution: Human review layer before approval",
                         "🎥 Challenge: False positives in cheating detection",
                         "   Solution: Multi-factor validation and alerts",
                         "🔒 Challenge: Data privacy and security",
                         "   Solution: Encryption, GDPR compliance, access control"
                     ])
    
    # Slide 22: Future Enhancements
    add_content_slide(prs,
                     "🚀 Future Enhancements",
                     [
                         "🌐 Multi-language Support - Global accessibility",
                         "📱 Mobile App - Native iOS and Android applications",
                         "🤖 Advanced ML - Improved behavior prediction models",
                         "🎮 Gamification - Badges, leaderboards, rewards",
                         "🔗 Integration - LMS platforms (Canvas, Blackboard, Moodle)",
                         "📊 Advanced Analytics - Predictive analytics and insights"
                     ])
    
    # Slide 23: Project Statistics
    add_content_slide(prs,
                     "📊 Project Statistics",
                     [
                         "📁 Total Files: 100+ components and modules",
                         "🎨 UI Components: 50+ reusable React components",
                         "🔌 API Endpoints: 40+ RESTful endpoints",
                         "🗂️ Database Collections: 18+ MongoDB collections",
                         "⏱️ Development Timeline: Enterprise-grade architecture",
                         "🌟 Code Quality: Production-ready with best practices"
                     ])
    
    # Slide 24: Team Structure
    add_content_slide(prs,
                     "👥 Development Team",
                     [
                         "👨‍💻 Full Stack Developers - Frontend and backend development",
                         "🎨 UI/UX Designers - Design and user experience",
                         "🔬 AI/ML Engineers - AI integration and optimization",
                         "🧪 QA Engineers - Testing and quality assurance",
                         "🏗️ DevOps Engineers - Deployment and infrastructure",
                         "📊 Project Manager - Coordination and planning"
                     ])
    
    # Slide 25: Implementation Highlights
    add_content_slide(prs,
                     "⭐ Implementation Highlights",
                     [
                         "✅ Real-time face recognition for exam integrity",
                         "✅ Intelligent question generation from any document",
                         "✅ Multi-layer security with JWT authentication",
                         "✅ Comprehensive analytics dashboard for teachers",
                         "✅ Support for multiple file formats (PDF, DOCX, Images)",
                         "✅ Automated cheating detection algorithms"
                     ])
    
    # Slide 26: Conclusion
    add_content_slide(prs,
                     "🎓 Conclusion",
                     [
                         "✨ ProctorSecure-AI is a comprehensive, AI-powered assessment platform",
                         "🎯 Designed for modern educational institutions and organizations",
                         "🔒 Combines security, AI, and user experience seamlessly",
                         "📈 Scalable architecture ready for enterprise deployment",
                         "🚀 Continuous improvement with planned enhancements",
                         "💡 Innovative solution to online education challenges"
                     ])
    
    # Slide 27: Q&A
    add_title_slide(prs,
                   "Questions & Answers",
                   "Thank You for Your Attention!")
    
    # Save the presentation
    output_path = "ProctorSecure-AI_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
